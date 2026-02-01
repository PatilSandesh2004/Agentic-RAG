# from pathlib import Path
# from typing import List

# import docx
# import pdfplumber
# from pptx import Presentation


# def load_document(file_path: str) -> str:
#     """
#     Load text from supported document types:
#     - TXT
#     - PDF
#     - DOCX
#     - PPTX
#     """

#     path = Path(file_path)
#     ext = path.suffix.lower()

#     if ext == ".txt":
#         return path.read_text(encoding="utf-8", errors="ignore")

#     elif ext == ".pdf":
#         return _load_pdf(path)

#     elif ext == ".docx":
#         return _load_docx(path)

#     elif ext == ".pptx":
#         return _load_pptx(path)

#     else:
#         raise ValueError(f"Unsupported file type: {ext}")


# # ---------------- PDF ----------------
# def _load_pdf(path: Path) -> str:
#     text = []
#     with pdfplumber.open(path) as pdf:
#         for page_num, page in enumerate(pdf.pages, start=1):
#             page_text = page.extract_text() or ""
#             if page_text.strip():
#                 text.append(f"\n--- PAGE {page_num} ---\n{page_text}")
#     return "\n\n".join(text)


# # ---------------- DOCX ----------------
# def _load_docx(path: Path) -> str:
#     doc = docx.Document(path)
#     return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())


# # ---------------- PPTX ----------------
# def _load_pptx(path: Path) -> str:
#     prs = Presentation(path)
#     text_runs: List[str] = []

#     for slide_index, slide in enumerate(prs.slides, start=1):
#         slide_text = []

#         for shape in slide.shapes:
#             if not shape.has_text_frame:
#                 continue

#             for paragraph in shape.text_frame.paragraphs:
#                 if paragraph.text.strip():
#                     slide_text.append(paragraph.text)

#         if slide_text:
#             text_runs.append(
#                 f"\n--- SLIDE {slide_index} ---\n" + "\n".join(slide_text)
#             )

#     return "\n\n".join(text_runs)


from pathlib import Path
from typing import List
import csv

import docx
import pdfplumber
import pandas as pd
from pptx import Presentation
from bs4 import BeautifulSoup


def load_document(file_path: str) -> str:
    """
    Load text from multiple document formats.
    """

    path = Path(file_path)
    ext = path.suffix.lower()

    if ext in [".txt", ".md"]:
        return path.read_text(encoding="utf-8", errors="ignore")

    elif ext == ".pdf":
        return _load_pdf(path)

    elif ext == ".docx":
        return _load_docx(path)

    elif ext == ".pptx":
        return _load_pptx(path)

    elif ext in [".xlsx", ".xls"]:
        return _load_excel(path)

    elif ext == ".csv":
        return _load_csv(path)

    elif ext in [".html", ".htm"]:
        return _load_html(path)

    else:
        raise ValueError(f"Unsupported file type: {ext}")


# ---------------- PDF ----------------
def _load_pdf(path: Path) -> str:
    text = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            page_text = page.extract_text() or ""
            if page_text.strip():
                text.append(f"\n--- PAGE {i} ---\n{page_text}")
    return "\n\n".join(text)


# ---------------- DOCX ----------------
def _load_docx(path: Path) -> str:
    doc = docx.Document(path)
    return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())


# ---------------- PPTX ----------------
def _load_pptx(path: Path) -> str:
    prs = Presentation(path)
    slides_text: List[str] = []

    for idx, slide in enumerate(prs.slides, start=1):
        slide_lines = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    if p.text.strip():
                        slide_lines.append(p.text)

        if slide_lines:
            slides_text.append(
                f"\n--- SLIDE {idx} ---\n" + "\n".join(slide_lines)
            )

    return "\n\n".join(slides_text)


# ---------------- EXCEL ----------------
def _load_excel(path: Path) -> str:
    sheets = pd.read_excel(path, sheet_name=None)
    text = []

    for sheet_name, df in sheets.items():
        rows = df.astype(str).values.tolist()
        row_text = [" | ".join(row) for row in rows if any(row)]
        if row_text:
            text.append(
                f"\n--- SHEET: {sheet_name} ---\n" + "\n".join(row_text)
            )

    return "\n\n".join(text)


# ---------------- CSV ----------------
def _load_csv(path: Path) -> str:
    rows = []
    with open(path, newline="", encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        for row in reader:
            if any(row):
                rows.append(" | ".join(row))
    return "\n".join(rows)


# ---------------- HTML ----------------
def _load_html(path: Path) -> str:
    soup = BeautifulSoup(
        path.read_text(encoding="utf-8", errors="ignore"),
        "lxml",
    )
    return soup.get_text(separator="\n")
